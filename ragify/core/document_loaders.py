from typing import List, Dict, Any, Optional, Union
import os
from ..config import get_config
import logging

logger = logging.getLogger("ragify.core.document_loaders")

# 适配langchain 1.0.5版本的导入
try:
    # 从langchain-core导入Document
    from langchain_core.documents import Document
    
    # 从langchain-community导入加载器（langchain 1.0.5将这些移到了community包）
    from langchain_community.document_loaders import (
        TextLoader, Docx2txtLoader, 
        DirectoryLoader, UnstructuredFileLoader
    )
    # 导入具体的PDF加载器（LangChain 1.0移除了通用PDFLoader）
    from langchain_community.document_loaders import PyPDFLoader as PDFLoader
    
    # 尝试导入UnstructuredImageLoader
    try:
        from langchain_community.document_loaders import UnstructuredImageLoader
    except ImportError:
        # 如果不存在，保持为None
        UnstructuredImageLoader = None
except ImportError as e:
    logger.warning(f"警告: 无法导入langchain组件: {e}")
    # 如果langchain完全无法导入，定义必要的类以避免崩溃
    class Document:
        def __init__(self, page_content, metadata=None):
            self.page_content = page_content
            self.metadata = metadata or {}
    
    class TextLoader:
        def __init__(self, file_path, encoding="utf-8"):
            self.file_path = file_path
            self.encoding = encoding
        def load(self):
            with open(self.file_path, 'r', encoding=self.encoding) as f:
                content = f.read()
            return [Document(page_content=content, metadata={"source": self.file_path})]
    
    # 为其他加载器定义基本的回退实现
    # PDFLoader使用PyPDF的简单实现作为回退
    class PDFLoader(TextLoader):
        pass
    
    class Docx2txtLoader(TextLoader):
        pass
    
    class DirectoryLoader:
        def __init__(self, *args, **kwargs):
            self.directory_path = args[0] if args else kwargs.get('directory_path', kwargs.get('path', '.'))
            self.glob = kwargs.get('glob', '**/*')
            self.loader_kwargs = kwargs.get('loader_kwargs', {})
            self.loaders = kwargs.get('loaders', {})
            self.show_progress = kwargs.get('show_progress', False)
            self.use_multithreading = kwargs.get('use_multithreading', True)
        def load(self):
            import glob
            files = glob.glob(os.path.join(self.directory_path, self.glob), recursive=True)
            files = [f for f in files if os.path.isfile(f)]
            documents = []
            for file in files:
                try:
                    ext = os.path.splitext(file)[1].lower()
                    if ext in self.loaders:
                        loader_cls, loader_kwargs = self.loaders[ext]
                        loader = loader_cls(file, **loader_kwargs)
                    elif ext == '.txt':
                        loader = TextLoader(file)
                    else:
                        continue  # 简化处理，只处理已知类型
                    documents.extend(loader.load())
                except Exception as e:
                    logger.error(f"无法加载文件 {file}: {e}")
            return documents
    
    class UnstructuredFileLoader(TextLoader):
        pass
    
    UnstructuredImageLoader = None


class MultiModalDocumentLoader:
    """
    多模态文档加载器，支持文本、PDF、Word、图像等多种格式
    """
    
    def __init__(self):
        self.config = get_config()
        self.multimodal_enabled = self.config.get("multimodal.enabled", True)
        self.image_enabled = self.config.get("multimodal.image_processor.enabled", True)
    
    def load_file(self, file_path: str) -> List[Document]:
        """
        加载单个文件
        """
        file_ext = os.path.splitext(file_path)[1].lower()

        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在: {file_path}")

        # 根据文件类型选择合适的加载器
        if file_ext == ".txt":
            loader = TextLoader(file_path, encoding="utf-8")
        elif file_ext == ".pdf":
            return self._load_pdf(file_path)
        elif file_ext == ".docx" or file_ext == ".doc":
            loader = Docx2txtLoader(file_path)
        elif file_ext == ".pptx":
            return self._load_pptx(file_path)
        elif file_ext in [".jpg", ".jpeg", ".png", ".gif"] and self.image_enabled and self.multimodal_enabled and UnstructuredImageLoader:
            loader = UnstructuredImageLoader(file_path)
        else:
            # 使用通用加载器作为后备
            loader = UnstructuredFileLoader(file_path)

        documents = loader.load()

        # 创建新的 Document 对象以避免修改原始对象
        result_docs = []
        for doc in documents:
            new_metadata = doc.metadata.copy()
            new_metadata["file_path"] = file_path
            new_metadata["file_type"] = file_ext
            new_metadata["source"] = file_path
            result_docs.append(Document(page_content=doc.page_content, metadata=new_metadata))

        return result_docs

    def _load_pdf(self, file_path: str) -> List[Document]:
        """Load PDF using pdftext (MinerU backend) with pypdf fallback."""
        content = None

        # Try pdftext first — better layout preservation
        try:
            from pdftext.extraction import plain_text_output
            content = plain_text_output(file_path, sort=True)
            if content and content.strip():
                logger.info(f"pdftext loaded: {file_path}")
        except Exception as e:
            logger.debug(f"pdftext failed, falling back to pypdf: {e}")

        # Fallback to pypdf
        if not content or not content.strip():
            try:
                loader = PDFLoader(file_path)
                docs = loader.load()
                content = "\n\n".join(d.page_content for d in docs)
                logger.info(f"pypdf loaded: {file_path}")
            except Exception as e:
                raise RuntimeError(f"无法解析 PDF 文件 {file_path}: {e}") from e

        doc = Document(
            page_content=content,
            metadata={
                "file_path": file_path,
                "file_type": ".pdf",
                "source": file_path,
            },
        )
        return [doc]

    def _load_pptx(self, file_path: str) -> List[Document]:
        """Load PPTX using python-pptx, extracting text from all slides."""
        from pptx import Presentation

        prs = Presentation(file_path)
        slides_text: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            lines: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            lines.append(text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            lines.append(row_text)
            if lines:
                slides_text.append(f"[Slide {i}]\n" + "\n".join(lines))

        if not slides_text:
            raise RuntimeError(f"PPTX 文件中未找到可提取的文本: {file_path}")

        content = "\n\n".join(slides_text)
        doc = Document(
            page_content=content,
            metadata={
                "file_path": file_path,
                "file_type": ".pptx",
                "source": file_path,
                "slide_count": len(prs.slides),
            },
        )
        return [doc]
    
    def load_directory(self, directory_path: str, glob_pattern: str = "**/*") -> List[Document]:
        """
        加载目录中的所有文件
        """
        if not os.path.isdir(directory_path):
            raise NotADirectoryError(f"目录不存在: {directory_path}")

        # 手动遍历文件并加载
        import glob
        files = glob.glob(os.path.join(directory_path, glob_pattern), recursive=True)
        files = [f for f in files if os.path.isfile(f)]

        documents = []
        for file_path in files:
            try:
                ext = os.path.splitext(file_path)[1].lower()
                if ext == '.txt':
                    loader = TextLoader(file_path, encoding="utf-8")
                elif ext == '.pdf':
                    docs = self._load_pdf(file_path)
                    for doc in docs:
                        documents.append(doc)
                    continue
                elif ext == '.docx' or ext == '.doc':
                    loader = Docx2txtLoader(file_path)
                elif ext == '.pptx':
                    docs = self._load_pptx(file_path)
                    for doc in docs:
                        documents.append(doc)
                    continue
                else:
                    continue  # 跳过不支持的文件类型
                docs = loader.load()
                for doc in docs:
                    new_metadata = doc.metadata.copy()
                    new_metadata["file_path"] = file_path
                    new_metadata["file_type"] = ext
                    documents.append(Document(page_content=doc.page_content, metadata=new_metadata))
            except Exception as e:
                logger.error(f"无法加载文件 {file_path}: {e}")

        return documents
    
    def load_from_config(self) -> List[Document]:
        """
        从配置中指定的数据目录加载文档
        """
        data_dir = self.config.get("base.data_dir", "./data")
        return self.load_directory(data_dir)


class ImageDocumentProcessor:
    """
    图像处理组件，专门处理图像文档
    """
    
    def __init__(self):
        self.config = get_config()
        self.max_size = self.config.get("multimodal.image_processor.max_size", 1024)
    
    def preprocess_image(self, image_path: str) -> Dict[str, Any]:
        """
        预处理图像文件
        """
        import PIL.Image
        from PIL import ImageOps
        
        # 打开并处理图像
        image = PIL.Image.open(image_path)
        
        # 调整图像大小
        image.thumbnail((self.max_size, self.max_size))
        
        # 转换为RGB模式（如果需要）
        if image.mode != "RGB":
            image = image.convert("RGB")
        
        # 返回图像信息
        return {
            "path": image_path,
            "width": image.width,
            "height": image.height,
            "mode": image.mode,
            "image": image
        }
    
    def extract_text_from_image(self, image_path: str) -> Optional[str]:
        """
        从图像中提取文本（OCR）
        注意：这需要额外安装pytesseract和Tesseract OCR
        """
        try:
            import pytesseract
            from PIL import Image
            
            image = Image.open(image_path)
            text = pytesseract.image_to_string(image, lang='chi_sim+eng')
            return text.strip()
        except ImportError:
            logger.warning("警告: pytesseract未安装，无法进行OCR提取")
            return None
        except Exception as e:
            logger.error(f"OCR处理错误: {e}")
            return None
